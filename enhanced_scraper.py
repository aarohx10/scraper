import asyncio
import json
import re
import tempfile
import os
import requests
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup
from pathlib import Path
from playwright.async_api import async_playwright
import fitz  # PyMuPDF
import argparse
import nltk
import os
nltk.data.path.append(os.getenv("NLTK_DATA", "/app/nltk_data"))
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.sentiment import SentimentIntensityAnalyzer
from langdetect import detect
from googletrans import Translator
import pytesseract
from PIL import Image
import io
import logging
from typing import Dict, List, Any, Optional
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import time

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants for timeouts and limits
MAX_EXECUTION_TIME = 8  # seconds (leaving 2s buffer for Vercel)
CHUNK_SIZE = 1  # Process one URL at a time
MAX_RETRIES = 3

class TimeoutError(Exception):
    pass

class ScrapingError(Exception):
    pass

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')
try:
    nltk.data.find('sentiment/vader_lexicon.zip')
except LookupError:
    nltk.download('vader_lexicon')

# Office document libraries
import docx
import openpyxl
from pptx import Presentation


async def search_google(query: str, num_results: int = 5) -> List[Dict[str, Any]]:
    """Search Google and return results."""
    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()
            
            # Navigate to Google
            await page.goto(f'https://www.google.com/search?q={query}')
            
            # Wait for results
            await page.wait_for_selector('div#search', timeout=30000)
            
            # Extract results
            results = await page.evaluate('''() => {
                const searchResults = [];
                const elements = document.querySelectorAll('div.g');
                elements.forEach((el, index) => {
                    if (index >= 5) return;
                    const titleEl = el.querySelector('h3');
                    const linkEl = el.querySelector('a');
                    const snippetEl = el.querySelector('div.VwiC3b');
                    
                    if (titleEl && linkEl && snippetEl) {
                        searchResults.push({
                            title: titleEl.innerText,
                            url: linkEl.href,
                            snippet: snippetEl.innerText
                        });
                    }
                });
                return searchResults;
            }''')
            
            await browser.close()
            return results
            
    except Exception as e:
        logger.error(f"Error in search_google: {str(e)}")
        return []


def is_valid_url(url):
    """Check if a URL is valid."""
    try:
        result = urlparse(url)
        return all([result.scheme, result.netloc])
    except:
        return False


def normalize_url(base_url, url):
    """Convert relative URLs to absolute URLs."""
    if is_valid_url(url):
        return url
    return urljoin(base_url, url)


def crawl_page(url):
    """Extract text and downloadable links from a webpage."""
    print(f"Crawling {url}")
    text = ""
    downloadables = []
    
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        response = requests.get(url, timeout=15, headers=headers)
        response.raise_for_status()
        
        content_type = response.headers.get('Content-Type', '').lower()
        
        # If the URL is a direct file
        if any(ct in content_type for ct in ['application/pdf', 'application/vnd.openxmlformats-officedocument', 'application/msword']):
            return "", [url]
            
        # If the URL is a direct text file
        if 'text/plain' in content_type:
            return response.text, []
        
        # Parse HTML content
        soup = BeautifulSoup(response.text, "html.parser")
        
        # Extract text content
        for script in soup(["script", "style"]):
            script.extract()
        text = soup.get_text(separator=" ", strip=True)
        
        # Extract links to downloadable files
        links = []
        for a in soup.find_all('a', href=True):
            href = a['href']
            full_url = normalize_url(url, href)
            if full_url and is_valid_url(full_url):
                links.append(full_url)
        
        # Filter for downloadable files
        downloadables = [
            link for link in links 
            if any(link.lower().endswith(ext) for ext in ['.pdf', '.txt', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'])
        ]
        
        return text, downloadables
        
    except Exception as e:
        print(f"Error crawling {url}: {e}")
        return "", []


def extract_pdf_text(url):
    """Download and extract text from a PDF file."""
    print(f"Extracting PDF from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            doc = fitz.open(temp_path)
            text = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text += page.get_text()
            doc.close()
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing PDF file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading PDF from {url}: {e}")
        return ""


def extract_docx_text(url):
    """Download and extract text from a DOCX file."""
    print(f"Extracting DOCX from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            doc = docx.Document(temp_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing DOCX file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading DOCX from {url}: {e}")
        return ""


def extract_xlsx_text(url):
    """Download and extract text from an XLSX file."""
    print(f"Extracting XLSX from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            workbook = openpyxl.load_workbook(temp_path, data_only=True)
            text = ""
            
            # Extract text from each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows():
                    row_text = " | ".join(str(cell.value) if cell.value is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing XLSX file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading XLSX from {url}: {e}")
        return ""


def extract_pptx_text(url):
    """Download and extract text from a PPTX file."""
    print(f"Extracting PPTX from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            presentation = Presentation(temp_path)
            text = ""
            
            for i, slide in enumerate(presentation.slides):
                text += f"Slide {i+1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                text += "\n"
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing PPTX file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading PPTX from {url}: {e}")
        return ""


def extract_txt_text(url):
    """Download and extract text from a text file."""
    print(f"Extracting TXT from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, timeout=15, headers=headers)
        r.raise_for_status()
        return r.text
    except Exception as e:
        print(f"Error extracting TXT from {url}: {e}")
        return ""


def clean_text(text):
    """Clean and normalize extracted text."""
    if not text:
        return ""
        
    # Remove URLs, emails, and special characters
    text = re.sub(r'https?://\S+|www\.\S+', '', text)
    text = re.sub(r'\S+@\S+\.\S+', '', text)
    text = re.sub(r'[^\w\s.,!?-]', ' ', text)
    
    # Remove extra whitespace and normalize
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    
    # Remove very long words (likely garbage)
    text = ' '.join(word for word in text.split() if len(word) < 50)
    
    return text


def save_json(data, filename="output.json"):
    """Save the results to a JSON file."""
    with open(filename, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Data saved to {filename}")


async def extract_document_text(url):
    """Extract text from document URLs based on file extension."""
    url_lower = url.lower()
    
    # Map file extensions to their extraction functions
    extractors = {
        '.pdf': extract_pdf_text,
        '.txt': extract_txt_text,
        '.docx': extract_docx_text,
        '.xlsx': extract_xlsx_text,
        '.pptx': extract_pptx_text
    }
    
    # Get the appropriate extractor based on file extension
    for ext, extractor in extractors.items():
        if url_lower.endswith(ext):
            try:
                return await asyncio.to_thread(extractor, url)
            except Exception as e:
                print(f"Error extracting {ext} from {url}: {e}")
                return ""
    
    # Handle legacy formats
    if any(url_lower.endswith(ext) for ext in ['.doc', '.xls', '.ppt']):
        print(f"Legacy format not supported: {url}")
    
    return ""


async def process_url(url):
    """Process a URL and extract all relevant information."""
    print(f"Processing {url}")
    result = {
        'url': url,
        'text': '',
        'language': 'unknown',
        'translated_text': '',
        'sentiment': {},
        'keywords': [],
        'downloadables': []
    }
    
    try:
        # Get webpage content
        text, downloadables = await asyncio.to_thread(crawl_page, url)
        result['downloadables'] = downloadables
        
        # Process main text
        if text:
            # Clean text first
            text = clean_text(text)
            
            # Detect language
            result['language'] = detect_language(text)
            
            # Translate if not English
            if result['language'] != 'en':
                result['translated_text'] = translate_text(text)
                processed_text = result['translated_text']
            else:
                processed_text = text
                
            result['text'] = processed_text
            
            # Analyze sentiment and extract keywords in parallel
            sentiment_task = asyncio.create_task(asyncio.to_thread(analyze_sentiment, processed_text))
            keywords_task = asyncio.create_task(asyncio.to_thread(extract_keywords, processed_text))
            
            result['sentiment'] = await sentiment_task
            result['keywords'] = await keywords_task
        
        # Process downloadable files in parallel
        if downloadables:
            doc_tasks = [extract_document_text(doc_url) for doc_url in downloadables]
            doc_texts = await asyncio.gather(*doc_tasks)
            
            for doc_text in doc_texts:
                if doc_text:
                    # Clean and process document text
                    doc_text = clean_text(doc_text)
                    
                    # Detect language and translate if needed
                    doc_lang = detect_language(doc_text)
                    if doc_lang != 'en':
                        doc_text = translate_text(doc_text)
                    
                    # Add to main text
                    result['text'] += "\n\n" + doc_text
                    
                    # Update sentiment and keywords in parallel
                    doc_sentiment_task = asyncio.create_task(asyncio.to_thread(analyze_sentiment, doc_text))
                    doc_keywords_task = asyncio.create_task(asyncio.to_thread(extract_keywords, doc_text))
                    
                    doc_sentiment = await doc_sentiment_task
                    doc_keywords = await doc_keywords_task
                    
                    # Update sentiment (weighted average)
                    result['sentiment'] = {
                        'compound': (result['sentiment'].get('compound', 0) + doc_sentiment['compound']) / 2,
                        'positive': (result['sentiment'].get('positive', 0) + doc_sentiment['positive']) / 2,
                        'negative': (result['sentiment'].get('negative', 0) + doc_sentiment['negative']) / 2,
                        'neutral': (result['sentiment'].get('neutral', 0) + doc_sentiment['neutral']) / 2
                    }
                    result['keywords'].extend(doc_keywords)
        
        # Remove duplicates and limit keywords
        result['keywords'] = list(set(result['keywords']))[:10]
        
    except Exception as e:
        print(f"Error processing {url}: {e}")
    
    return result


@retry(
    stop=stop_after_attempt(MAX_RETRIES),
    wait=wait_exponential(multiplier=1, min=4, max=10),
    retry=retry_if_exception_type((ScrapingError, TimeoutError))
)
async def process_url_with_timeout(url: str) -> Dict[str, Any]:
    """Process a single URL with timeout."""
    start_time = time.time()
    
    try:
        result = await asyncio.wait_for(
            process_url(url),
            timeout=MAX_EXECUTION_TIME
        )
        
        if time.time() - start_time > MAX_EXECUTION_TIME:
            raise TimeoutError(f"Processing {url} took too long")
            
        return result
    except asyncio.TimeoutError:
        raise TimeoutError(f"Timeout while processing {url}")
    except Exception as e:
        raise ScrapingError(f"Error processing {url}: {str(e)}")


async def main(query: str, num_results: int = 5, continuation_token: str = None) -> Dict[str, Any]:
    """Main function to process search request with strict timeouts and retries."""
    start_time = time.time()
    
    try:
        # If we have a continuation token, resume processing
        if continuation_token:
            try:
                token_data = json.loads(continuation_token)
                processed_urls = set(token_data.get('processed_urls', []))
                remaining_urls = token_data.get('remaining_urls', [])
                results = token_data.get('results', [])
            except json.JSONDecodeError:
                return {
                    'status': 'error',
                    'error': 'Invalid continuation token'
                }
        else:
            # Start new processing
            processed_urls = set()
            results = []
            try:
                search_results = await asyncio.wait_for(
                    search_google(query, num_results),
                    timeout=MAX_EXECUTION_TIME
                )
                remaining_urls = [result['url'] for result in search_results]
                results.extend(search_results)
            except asyncio.TimeoutError:
                return {
                    'status': 'error',
                    'error': 'Timeout during initial search'
                }

        # Process URLs one at a time with strict timeouts
        while remaining_urls and len(processed_urls) < num_results:
            # Check if we're approaching the time limit
            if time.time() - start_time > MAX_EXECUTION_TIME:
                token_data = {
                    'processed_urls': list(processed_urls),
                    'remaining_urls': remaining_urls,
                    'results': results
                }
                return {
                    'status': 'in_progress',
                    'continuation_token': json.dumps(token_data),
                    'results': results,
                    'time_elapsed': time.time() - start_time
                }

            # Process next URL
            current_url = remaining_urls[0]
            remaining_urls = remaining_urls[1:]
            
            try:
                result = await process_url_with_timeout(current_url)
                processed_urls.add(current_url)
                results.append(result)
            except (TimeoutError, ScrapingError) as e:
                logger.error(f"Error processing {current_url}: {str(e)}")
                # Don't add failed URLs back to the queue
                continue

        return {
            'status': 'success',
            'query': query,
            'results': results,
            'time_elapsed': time.time() - start_time
        }
        
    except Exception as e:
        logger.error(f"Unexpected error in main function: {str(e)}")
        return {
            'status': 'error',
            'error': str(e),
            'time_elapsed': time.time() - start_time
        }


def detect_language(text):
    """Detect the language of the given text."""
    try:
        return detect(text)
    except:
        return 'unknown'


def translate_text(text, target_lang='en'):
    """Translate text to target language."""
    try:
        translator = Translator()
        result = translator.translate(text, dest=target_lang)
        return result.text
    except Exception as e:
        print(f"Translation error: {e}")
        return text


def extract_text_from_image(image_data):
    """Extract text from image using OCR."""
    try:
        image = Image.open(io.BytesIO(image_data))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"OCR error: {e}")
        return ""


def analyze_sentiment(text):
    """Analyze sentiment of text using VADER."""
    try:
        sia = SentimentIntensityAnalyzer()
        scores = sia.polarity_scores(text)
        return {
            'compound': scores['compound'],
            'positive': scores['pos'],
            'negative': scores['neg'],
            'neutral': scores['neu']
        }
    except Exception as e:
        print(f"Sentiment analysis error: {e}")
        return {'compound': 0, 'positive': 0, 'negative': 0, 'neutral': 0}


def extract_keywords(text, num_keywords=10):
    """Extract key phrases from text."""
    try:
        # Tokenize and remove stopwords
        tokens = word_tokenize(text.lower())
        stop_words = set(stopwords.words('english'))
        filtered_tokens = [word for word in tokens if word.isalnum() and word not in stop_words]
        
        # Get word frequency
        freq_dist = nltk.FreqDist(filtered_tokens)
        
        # Return top keywords
        return [word for word, freq in freq_dist.most_common(num_keywords)]
    except Exception as e:
        print(f"Keyword extraction error: {e}")
        return []


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Enhanced Office Document Scraper")
    parser.add_argument("query", help="Search query")
    parser.add_argument("--results", type=int, default=5, help="Number of results to process")
    args = parser.parse_args()
    
    try:
        asyncio.run(main(args.query, args.results))
    except KeyboardInterrupt:
        print("\nScraping interrupted by user")
    except Exception as e:
        print(f"Error running scraper: {e}") 